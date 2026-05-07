#!/usr/bin/env python3
"""
Document-RAG 开题展示 PPT 生成脚本 (8页版)
使用 python-pptx 生成可编辑 PPTX
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from lxml import etree
import os

# ============== 配置 ==============
OUTPUT_PATH = "/Users/lianglihang/Downloads/Document-RAG/文档理解与多模态GraphRAG检索问答系统-开题展示.pptx"
ARCH_IMAGE = "/tmp/pptx_extract/ppt/media/image1.png"
DEMO_IMAGE = "/tmp/pptx_extract/ppt/media/image2.png"

# 16:9 尺寸 (1280x720)
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# ============== 颜色 ==============
BG = RGBColor(0xF3, 0xF2, 0xEE)
SURFACE = RGBColor(0xFF, 0xFE, 0xFC)
SURFACE_ALT = RGBColor(0xF4, 0xF8, 0xF6)
INK = RGBColor(0x14, 0x2A, 0x2B)
BODY = RGBColor(0x58, 0x6D, 0x70)
MUTE = RGBColor(0x7A, 0x8D, 0x90)
ACCENT = RGBColor(0x0E, 0x7D, 0x6A)
ACCENT_SOFT = RGBColor(0xE6, 0xF1, 0xEE)
ACCENT_DEEP = RGBColor(0x0A, 0x5C, 0x50)
BORDER = RGBColor(0xD7, 0xE1, 0xDE)
WARM = RGBColor(0xF4, 0xEF, 0xE6)
PEACH = RGBColor(0xEF, 0xE6, 0xD9)

# ============== 辅助函数 ==============
def set_slide_size(prs):
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

def add_background(slide, variant="default"):
    """添加背景形状"""
    # 主背景
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    bg.line.fill.background()
    
    # 顶部条
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, Inches(0.917))
    band.fill.solid()
    band.fill.fore_color.rgb = RGBColor(0xEC, 0xEF, 0xEB)
    band.line.fill.background()
    
    # 右上角装饰块
    accent_block = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(10.208), 0, Inches(3.125), Inches(0.917))
    accent_block.fill.solid()
    accent_block.fill.fore_color.rgb = ACCENT_SOFT
    accent_block.line.fill.background()
    
    # 左下角装饰块
    if variant == "warm":
        footer_block = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(6.403), Inches(2.583), Inches(1.097))
        footer_block.fill.solid()
        footer_block.fill.fore_color.rgb = PEACH
    else:
        footer_block = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(6.403), Inches(2.583), Inches(1.097))
        footer_block.fill.solid()
        footer_block.fill.fore_color.rgb = RGBColor(0xEA, 0xF3, 0xF0)
    footer_block.line.fill.background()
    
    # 画布
    canvas = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.417), Inches(0.417), Inches(12.5), Inches(6.667))
    canvas.fill.solid()
    canvas.fill.fore_color.rgb = SURFACE
    canvas.line.color.rgb = BORDER
    canvas.line.width = Pt(1)
    
    # 画布左侧装饰条
    spine = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.417), Inches(0.417), Inches(0.104), Inches(6.667))
    spine.fill.solid()
    if variant == "warm":
        spine.fill.fore_color.rgb = WARM
    else:
        spine.fill.fore_color.rgb = SURFACE_ALT
    spine.line.fill.background()

def add_textbox(slide, left, top, width, height, text, font_size=18, bold=False, color=INK, align=PP_ALIGN.LEFT, font_name="PingFang SC"):
    """添加文本框"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return txBox

def add_header(slide, slide_no, kicker, total=8):
    """添加页眉"""
    # 左侧装饰条
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.812), Inches(0.792), Inches(0.042), Inches(0.25))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()
    
    # kicker
    add_textbox(slide, Inches(0.958), Inches(0.75), Inches(6.0), Inches(0.3), kicker,
                font_size=12, bold=True, color=ACCENT_DEEP, font_name="Aptos")
    
    # 页码 pill
    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(11.708), Inches(0.681), Inches(1.042), Inches(0.472))
    pill.fill.solid()
    pill.fill.fore_color.rgb = SURFACE_ALT
    pill.line.color.rgb = BORDER
    pill.line.width = Pt(1)
    
    add_textbox(slide, Inches(11.875), Inches(0.778), Inches(0.75), Inches(0.25),
                f"{slide_no:02d} / {total:02d}", font_size=12, bold=True, color=ACCENT_DEEP, align=PP_ALIGN.CENTER, font_name="Aptos")
    
    # 分隔线
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.812), Inches(1.306), Inches(11.771), Inches(0.021))
    line.fill.solid()
    line.fill.fore_color.rgb = BORDER
    line.line.fill.background()

def add_panel(slide, left, top, width, height, fill=SURFACE, line_color=BORDER, accent_color=None):
    """添加面板"""
    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    panel.fill.solid()
    panel.fill.fore_color.rgb = fill
    panel.line.color.rgb = line_color
    panel.line.width = Pt(1)
    
    if accent_color:
        accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(0.062), height)
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = accent_color
        accent_bar.line.fill.background()
    
    return panel

def add_bullet_list(slide, left, top, width, height, items, font_size=16, color=BODY):
    """添加项目符号列表"""
    text = "\n".join([f"• {item}" for item in items])
    return add_textbox(slide, left, top, width, height, text, font_size=font_size, color=color)

# ============== 各页生成函数 ==============
def slide_cover(slide):
    """第1页：开始页"""
    add_background(slide, "warm")
    
    # 左侧装饰条
    rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.854), Inches(1.639), Inches(0.062), Inches(6.083))
    rule.fill.solid()
    rule.fill.fore_color.rgb = ACCENT
    rule.line.fill.background()
    
    # 课程标签
    add_textbox(slide, Inches(1.104), Inches(1.611), Inches(3.0), Inches(0.25),
                "课程作业开题答辩", font_size=12, bold=True, color=ACCENT_DEEP, font_name="Aptos")
    
    # 课程信息
    add_textbox(slide, Inches(1.104), Inches(2.056), Inches(5.0), Inches(0.278),
                "多模态大模型原理与应用 · 题目 3", font_size=12, bold=True, color=MUTE, font_name="Aptos")
    
    # 标题
    add_textbox(slide, Inches(1.042), Inches(2.611), Inches(9.0), Inches(1.889),
                "文档理解与多模态 GraphRAG 检索问答系统",
                font_size=44, bold=True, color=INK)
    
    # 副标题
    add_textbox(slide, Inches(1.104), Inches(4.694), Inches(8.0), Inches(0.917),
                "基于 qwen3-vl-8b-instruct、FAISS 与 networkx 的可追溯文档智能助手",
                font_size=18, color=BODY)
    
    # 右侧摘要面板
    add_panel(slide, Inches(10.417), Inches(1.806), Inches(2.833), Inches(2.417),
              fill=SURFACE_ALT, accent_color=ACCENT)
    add_textbox(slide, Inches(10.833), Inches(2.167), Inches(2.0), Inches(0.25),
                "核心技术栈", font_size=13, bold=True, color=ACCENT_DEEP, font_name="Aptos")
    add_bullet_list(slide, Inches(10.833), Inches(2.694), Inches(2.417), Inches(1.111),
                   ["qwen3-vl-8b-instruct 文档解析", "FAISS 向量索引 + networkx 图谱", "DashScope API 多模态问答"],
                   font_size=16, color=BODY)
    
    # 底部标签
    tags = ["多模态 GraphRAG", "DocVQA / ChartQA", "qwen3-vl-8b-instruct", "FAISS", "networkx"]
    tag_x = Inches(1.083)
    for tag in tags:
        width = max(Inches(1.2), min(Inches(2.5), len(tag) * Inches(0.18)))
        tag_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, tag_x, Inches(6.111), width, Inches(0.444))
        tag_shape.fill.solid()
        tag_shape.fill.fore_color.rgb = SURFACE
        tag_shape.line.color.rgb = BORDER
        tag_shape.line.width = Pt(1)
        add_textbox(slide, tag_x + Inches(0.146), Inches(6.222), width - Inches(0.292), Inches(0.222),
                   tag, font_size=12, bold=True, color=ACCENT_DEEP, align=PP_ALIGN.CENTER)
        tag_x += width + Inches(0.125)
    
    # 底部卡片
    add_panel(slide, Inches(1.083), Inches(7.194), Inches(3.208), Inches(0.833), accent_color=ACCENT)
    add_textbox(slide, Inches(1.333), Inches(7.361), Inches(1.5), Inches(0.25),
                "核心技术栈", font_size=12, bold=True, color=MUTE)
    add_textbox(slide, Inches(1.333), Inches(7.611), Inches(2.8), Inches(0.417),
                "qwen3-vl-8b-instruct + FAISS + networkx + DashScope", font_size=15, bold=True, color=INK)
    
    add_panel(slide, Inches(4.479), Inches(7.194), Inches(3.208), Inches(0.833), accent_color=ACCENT)
    add_textbox(slide, Inches(4.729), Inches(7.361), Inches(1.5), Inches(0.25),
                "当前状态", font_size=12, bold=True, color=MUTE)
    add_textbox(slide, Inches(4.729), Inches(7.611), Inches(2.8), Inches(0.417),
                "原型系统已完成，进入评测阶段", font_size=15, bold=True, color=INK)

def slide_background(slide):
    """第2页：任务背景"""
    add_background(slide)
    add_header(slide, 2, "01 · 任务背景")
    
    # 标题
    add_textbox(slide, Inches(0.917), Inches(1.528), Inches(10.0), Inches(0.833),
                "为什么纯文本 RAG 不够", font_size=36, bold=True, color=INK)
    
    # 三个卡片
    cards = [
        ("A", "视觉信息丢失", "图表趋势、版面结构、颜色标记、坐标轴和表格边界很难被纯文本准确保留。"),
        ("B", "结构关系不足", "同一页内的图文对应、跨页段落延续、图号与结论之间的支撑关系无法仅靠向量相似度稳定表达。"),
        ("C", "证据链不可解释", "向量检索可以找到相似片段，但不擅长回答\"某结论由哪些图表和实验支撑\"这类多跳问题。"),
    ]
    
    card_w = Inches(3.646)
    gap = Inches(0.25)
    for i, (idx, title, body) in enumerate(cards):
        x = Inches(0.792) + i * (card_w + gap)
        y = Inches(3.528)
        add_panel(slide, x, y, card_w, Inches(3.472), accent_color=ACCENT)
        
        # 索引圆圈
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.25), y + Inches(0.229), Inches(0.472), Inches(0.472))
        circle.fill.solid()
        circle.fill.fore_color.rgb = ACCENT_SOFT
        circle.line.fill.background()
        add_textbox(slide, x + Inches(0.403), y + Inches(0.354), Inches(0.167), Inches(0.167),
                   idx, font_size=12, bold=True, color=ACCENT_DEEP, align=PP_ALIGN.CENTER)
        
        add_textbox(slide, x + Inches(0.25), y + Inches(1.0), Inches(3.0), Inches(0.417),
                   title, font_size=22, bold=True, color=INK)
        add_textbox(slide, x + Inches(0.25), y + Inches(1.694), Inches(3.104), Inches(1.278),
                   body, font_size=15, color=BODY)
    
    # 底部条
    add_panel(slide, Inches(0.792), Inches(7.5), Inches(11.667), Inches(0.964), fill=SURFACE_ALT)
    add_textbox(slide, Inches(1.0), Inches(7.694), Inches(1.5), Inches(0.25),
                "系统要做什么", font_size=13, bold=True, color=ACCENT_DEEP)
    add_textbox(slide, Inches(1.0), Inches(7.944), Inches(11.0), Inches(0.417),
                "上传多页 PDF 或图像后，系统同时构建文本/图表片段、页面图像证据、关键词/关系节点和跨页连接，使回答既能利用多模态证据，又能给出可追溯的页码、节点和关系依据。",
                font_size=16, bold=True, color=INK)

def slide_architecture(slide):
    """第3页：架构图"""
    add_background(slide)
    add_header(slide, 3, "02 · 系统总览")
    
    # 标题
    add_textbox(slide, Inches(0.917), Inches(1.528), Inches(10.0), Inches(0.833),
                "系统怎么跑起来", font_size=36, bold=True, color=INK)
    add_textbox(slide, Inches(0.917), Inches(2.361), Inches(10.0), Inches(0.417),
                "知识库构建 + 在线问答 + 前端展示三条链路", font_size=16, color=BODY)
    
    # 图片
    if os.path.exists(ARCH_IMAGE):
        slide.shapes.add_picture(ARCH_IMAGE, Inches(0.75), Inches(2.944), width=Inches(11.5))
    else:
        add_textbox(slide, Inches(0.75), Inches(2.944), Inches(11.5), Inches(5.833),
                   "[系统架构图]", font_size=24, color=MUTE, align=PP_ALIGN.CENTER)
    
    # 底部说明
    add_panel(slide, Inches(0.75), Inches(9.028), Inches(11.5), Inches(0.417), fill=SURFACE_ALT)
    add_textbox(slide, Inches(1.0), Inches(9.111), Inches(11.0), Inches(0.278),
                "PyMuPDF → qwen3-vl-8b-instruct → FAISS + networkx → 多模态问答",
                font_size=14, color=BODY, align=PP_ALIGN.CENTER)

def slide_experiment(slide):
    """第4页：实验设计"""
    add_background(slide)
    add_header(slide, 4, "03 · 实验设计")
    
    # 标题
    add_textbox(slide, Inches(0.917), Inches(1.528), Inches(10.0), Inches(0.833),
                "评测什么、怎么评", font_size=36, bold=True, color=INK)
    
    # 数据集卡片
    datasets = [
        ("数据集 A", "自建长文档测试集", "3-5 篇含图表、表格、跨页段落和结论引用的论文，人工编写问题。"),
        ("数据集 B", "DocVQA 子集", "公开 DocVQA 验证样本，评估单页文档问答抽取式准确性。"),
        ("数据集 C", "ChartQA 子集", "图表推理问题，观察页面原图和图表描述的作用。"),
    ]
    
    ds_w = Inches(3.646)
    for i, (tag, title, body) in enumerate(datasets):
        x = Inches(0.792) + i * (ds_w + Inches(0.229))
        y = Inches(3.444)
        add_panel(slide, x, y, ds_w, Inches(2.028))
        add_textbox(slide, x + Inches(0.25), y + Inches(0.306), Inches(1.2), Inches(0.222),
                   tag, font_size=12, bold=True, color=ACCENT_DEEP)
        add_textbox(slide, x + Inches(0.25), y + Inches(0.667), Inches(2.396), Inches(0.333),
                   title, font_size=22, bold=True, color=INK)
        add_textbox(slide, x + Inches(0.25), y + Inches(1.194), Inches(3.0), Inches(0.583),
                   body, font_size=14, color=BODY)
    
    # 消融实验
    add_panel(slide, Inches(0.792), Inches(5.972), Inches(11.667), Inches(2.639),
              fill=SURFACE_ALT, accent_color=ACCENT)
    add_textbox(slide, Inches(1.042), Inches(6.222), Inches(1.5), Inches(0.222),
                "消融实验", font_size=12, bold=True, color=ACCENT_DEEP)
    add_textbox(slide, Inches(1.042), Inches(6.528), Inches(5.0), Inches(0.417),
                "Full / No-Graph / No-Image / Short-Answer", font_size=20, bold=True, color=INK)
    add_bullet_list(slide, Inches(1.042), Inches(7.028), Inches(11.0), Inches(1.5),
                   ["Full：向量召回 + 图谱扩展 + 社区 profile + 页面图像",
                    "No-Graph：关闭图谱扩展，观察多跳问题下降幅度",
                    "No-Image：关闭页面原图输入，观察图表题下降幅度",
                    "Short-Answer Eval：约束输出为短答案，计算 ANLS/EM"],
                   font_size=15, color=BODY)

def slide_demo(slide):
    """第5页：Demo"""
    add_background(slide)
    add_header(slide, 5, "04 · Demo")
    
    # 标题
    add_textbox(slide, Inches(0.917), Inches(1.528), Inches(10.0), Inches(0.833),
                "Knowledge Flow Studio", font_size=36, bold=True, color=INK)
    add_textbox(slide, Inches(0.917), Inches(2.361), Inches(10.0), Inches(0.417),
                "问答、知识图谱、节点检查、页面预览四块功能", font_size=16, color=BODY)
    
    # 图片
    if os.path.exists(DEMO_IMAGE):
        slide.shapes.add_picture(DEMO_IMAGE, Inches(0.75), Inches(2.778), width=Inches(11.5))
    else:
        add_textbox(slide, Inches(0.75), Inches(2.778), Inches(11.5), Inches(6.111),
                   "[Demo 截图]", font_size=24, color=MUTE, align=PP_ALIGN.CENTER)
    
    # 功能标签
    features = ["问答", "知识图谱", "节点检查", "页面预览"]
    feat_w = Inches(2.0)
    feat_gap = Inches(0.417)
    start_x = Inches(0.75) + (Inches(11.5) - (len(features) * feat_w + (len(features) - 1) * feat_gap)) / 2
    
    for i, feat in enumerate(features):
        fx = start_x + i * (feat_w + feat_gap)
        tag = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, fx, Inches(9.167), feat_w, Inches(0.389))
        tag.fill.solid()
        tag.fill.fore_color.rgb = ACCENT_SOFT
        tag.line.color.rgb = BORDER
        add_textbox(slide, fx, Inches(9.222), feat_w, Inches(0.278),
                   feat, font_size=13, bold=True, color=ACCENT_DEEP, align=PP_ALIGN.CENTER)

def slide_budget(slide):
    """第6页：算力预算与时间计划"""
    add_background(slide, "warm")
    add_header(slide, 6, "05 · 算力和计划")
    
    # 标题
    add_textbox(slide, Inches(0.917), Inches(1.528), Inches(10.0), Inches(0.833),
                "算力预算与时间计划", font_size=36, bold=True, color=INK)
    
    # 左侧：资源
    add_panel(slide, Inches(0.792), Inches(3.444), Inches(5.625), Inches(5.278), accent_color=ACCENT)
    add_textbox(slide, Inches(1.083), Inches(3.75), Inches(2.0), Inches(0.278),
                "算力与成本", font_size=15, bold=True, color=ACCENT_DEEP)
    
    resources = [
        ("本地计算", "CPU 即可", "PDF 渲染、FAISS、networkx 和 Flask 均可本地运行"),
        ("GPU", "非必需", "使用 DashScope API 调用 qwen3-vl-8b-instruct"),
        ("存储", "百 MB 到数 GB", "页面图片、FAISS 索引、JSON 解析结果和图谱文件"),
        ("API 成本", "与页数和问题数相关", "先小样本评测，再扩展；复用 build meta 避免重复入库"),
    ]
    
    ry = Inches(4.306)
    for name, value, detail in resources:
        add_textbox(slide, Inches(1.083), ry, Inches(1.5), Inches(0.25),
                   name, font_size=14, bold=True, color=ACCENT_DEEP)
        add_textbox(slide, Inches(2.188), ry, Inches(1.5), Inches(0.25),
                   value, font_size=14, bold=True, color=INK)
        add_textbox(slide, Inches(1.083), ry + Inches(0.278), Inches(5.0), Inches(0.417),
                   detail, font_size=13, color=BODY)
        ry += Inches(0.861)
    
    # 右上：进度
    add_panel(slide, Inches(6.667), Inches(3.444), Inches(5.833), Inches(2.5),
              fill=SURFACE_ALT, accent_color=ACCENT)
    add_textbox(slide, Inches(6.917), Inches(3.75), Inches(2.0), Inches(0.278),
                "当前进度", font_size=15, bold=True, color=ACCENT_DEEP)
    
    progress = [
        "工程骨架：已完成",
        "文档解析：已完成初版",
        "向量索引：已完成",
        "图谱构建：已完成初版",
        "GraphRAG 问答：已完成初版",
        "Web 演示：已完成初版",
        "评测脚本：已完成框架",
    ]
    add_bullet_list(slide, Inches(6.917), Inches(4.111), Inches(5.333), Inches(1.667),
                   progress, font_size=13, color=BODY)
    
    # 右下：计划
    add_panel(slide, Inches(6.667), Inches(6.222), Inches(5.833), Inches(2.5), accent_color=ACCENT)
    add_textbox(slide, Inches(6.917), Inches(6.528), Inches(2.5), Inches(0.278),
                "后续 4 周计划", font_size=15, bold=True, color=ACCENT_DEEP)
    
    plan = [
        "第 1 周：清理图谱噪声、优化关键词和关系过滤",
        "第 2 周：构建自建长文档测试集，补充 DocVQA 小样本",
        "第 3 周：跑 Full/No-Graph/No-Image 消融实验",
        "第 4 周：完善前端演示、整理最终报告和答辩材料",
    ]
    add_bullet_list(slide, Inches(6.917), Inches(6.889), Inches(5.333), Inches(1.667),
                   plan, font_size=13, color=BODY)

def slide_references(slide):
    """第7页：参考文献"""
    add_background(slide)
    add_header(slide, 7, "06 · 参考文献")
    
    # 标题
    add_textbox(slide, Inches(0.917), Inches(1.528), Inches(10.0), Inches(0.833),
                "核心来源", font_size=36, bold=True, color=INK)
    
    # 参考文献面板
    add_panel(slide, Inches(0.792), Inches(3.444), Inches(11.667), Inches(5.278), accent_color=ACCENT)
    
    refs = [
        "[1] Lewis, P. et al. (2020). Retrieval-Augmented Generation for Knowledge-Intensive NLP Tasks. NeurIPS 2020.",
        "[2] Edge, D. et al. (2024). From Local to Global: A Graph RAG Approach to Query-Focused Summarization. arXiv:2404.16130.",
        "[3] Microsoft GraphRAG Documentation. Query Engine Overview and Indexing Overview.",
        "[4] Mathew, M. et al. (2021). DocVQA: A Dataset for VQA on Document Images. WACV 2021.",
        "[5] Masry, A. et al. (2022). ChartQA: A Benchmark for Question Answering about Charts. Findings of ACL 2022.",
        "[6] Faysse, M. et al. (2024). ColPali: Efficient Document Retrieval with Vision Language Models. arXiv:2407.01449.",
        "[7] Qwen Team. Qwen3-VL official repository and Qwen3-VL-8B-Instruct model card.",
        "[8] Dong, K. et al. (2025). Benchmarking Retrieval-Augmented Multimodal Generation for Document Question Answering. MMDocRAG.",
        "[9] Bu, C. et al. (2025). Query-Driven Multimodal GraphRAG: Dynamic Local Knowledge Graph Construction for Online Reasoning. Findings of ACL 2025.",
    ]
    
    refs_text = "\n\n".join(refs)
    add_textbox(slide, Inches(1.083), Inches(3.75), Inches(11.0), Inches(4.833),
                refs_text, font_size=13, color=BODY)

def slide_closing(slide):
    """第8页：结束"""
    add_background(slide, "warm")
    
    # 左侧装饰条
    rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.854), Inches(1.639), Inches(0.062), Inches(6.083))
    rule.fill.solid()
    rule.fill.fore_color.rgb = ACCENT
    rule.line.fill.background()
    
    # kicker
    add_textbox(slide, Inches(1.104), Inches(2.056), Inches(5.0), Inches(0.278),
                "07 · 结束", font_size=12, bold=True, color=MUTE, font_name="Aptos")
    
    # 标题
    add_textbox(slide, Inches(1.042), Inches(2.611), Inches(8.0), Inches(0.833),
                "谢谢", font_size=56, bold=True, color=INK)
    
    # 副标题
    add_textbox(slide, Inches(1.104), Inches(3.611), Inches(9.0), Inches(0.556),
                "多模态 GraphRAG · 可追溯文档问答 · qwen3-vl-8b-instruct + FAISS + networkx",
                font_size=18, color=BODY)
    
    # 收束语
    add_panel(slide, Inches(0.792), Inches(4.722), Inches(11.667), Inches(1.389), fill=SURFACE_ALT)
    add_textbox(slide, Inches(1.083), Inches(5.0), Inches(11.0), Inches(0.833),
                "从入库开始同时构建向量索引和文档图谱，使文本、图表、页面和跨页关系处于同一证据空间。",
                font_size=17, bold=True, color=INK)
    
    # 关键词标签
    keywords = ["GraphRAG", "FAISS", "networkx", "qwen3-vl-8b-instruct", "DashScope", "PyMuPDF", "Flask"]
    kw_x = Inches(1.083)
    for kw in keywords:
        width = max(Inches(1.0), len(kw) * Inches(0.16))
        tag = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, kw_x, Inches(6.667), width, Inches(0.389))
        tag.fill.solid()
        tag.fill.fore_color.rgb = SURFACE
        tag.line.color.rgb = BORDER
        add_textbox(slide, kw_x + Inches(0.083), Inches(6.722), width - Inches(0.167), Inches(0.278),
                   kw, font_size=12, bold=True, color=ACCENT_DEEP, align=PP_ALIGN.CENTER)
        kw_x += width + Inches(0.146)

# ============== 主程序 ==============
def main():
    prs = Presentation()
    set_slide_size(prs)
    
    # 清空默认幻灯片
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]
    
    # 添加8页
    slide_cover(prs.slides.add_slide(prs.slide_layouts[6]))  # 空白布局
    slide_background(prs.slides.add_slide(prs.slide_layouts[6]))
    slide_architecture(prs.slides.add_slide(prs.slide_layouts[6]))
    slide_experiment(prs.slides.add_slide(prs.slide_layouts[6]))
    slide_demo(prs.slides.add_slide(prs.slide_layouts[6]))
    slide_budget(prs.slides.add_slide(prs.slide_layouts[6]))
    slide_references(prs.slides.add_slide(prs.slide_layouts[6]))
    slide_closing(prs.slides.add_slide(prs.slide_layouts[6]))
    
    # 保存
    prs.save(OUTPUT_PATH)
    print(f"PPT saved to: {OUTPUT_PATH}")
    print(f"Total slides: {len(prs.slides)}")

if __name__ == "__main__":
    main()
